#!/usr/bin/env python3
"""
Generate PDF and PowerPoint presentations from documentation
"""

import os
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image, Table, TableStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

class PDFGenerator:
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self.title_style = ParagraphStyle(
            'CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a1a1a'),
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        self.heading_style = ParagraphStyle(
            'CustomHeading',
            parent=self.styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#2c3e50'),
            spaceAfter=12,
            spaceBefore=12,
            fontName='Helvetica-Bold'
        )
        self.body_style = ParagraphStyle(
            'CustomBody',
            parent=self.styles['BodyText'],
            fontSize=11,
            alignment=TA_JUSTIFY,
            spaceAfter=10
        )
    
    def generate_comprehensive_pdf(self, output_file):
        """Generate comprehensive PDF documentation"""
        doc = SimpleDocTemplate(output_file, pagesize=letter,
                                rightMargin=72, leftMargin=72,
                                topMargin=72, bottomMargin=18)
        
        story = []
        
        # Title Page
        story.append(Spacer(1, 2*inch))
        story.append(Paragraph("RosterBhai", self.title_style))
        story.append(Spacer(1, 0.2*inch))
        story.append(Paragraph("System Documentation", self.title_style))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"Version 2.0.0<br/>{datetime.now().strftime('%B %d, %Y')}", 
                              ParagraphStyle('center', alignment=TA_CENTER, fontSize=14)))
        story.append(PageBreak())
        
        # Table of Contents
        story.append(Paragraph("Table of Contents", self.heading_style))
        story.append(Spacer(1, 0.2*inch))
        
        toc_items = [
            "1. Executive Summary",
            "2. System Architecture Overview",
            "3. Multi-Tenant Architecture",
            "4. System Components",
            "5. Data Models",
            "6. Authentication & Authorization",
            "7. API Architecture",
            "8. Security Features",
            "9. Diagrams",
        ]
        
        for item in toc_items:
            story.append(Paragraph(item, self.body_style))
        
        story.append(PageBreak())
        
        # Executive Summary
        story.append(Paragraph("1. Executive Summary", self.heading_style))
        story.append(Paragraph(
            "RosterBhai is a modern, multi-tenant SaaS application designed for roster and shift management. "
            "Built using Next.js 14 with TypeScript, it provides a comprehensive solution for organizations "
            "managing employee shifts across different teams.",
            self.body_style
        ))
        story.append(Spacer(1, 0.2*inch))
        
        # Key Features
        story.append(Paragraph("Key Features:", 
                              ParagraphStyle('bold', fontName='Helvetica-Bold', fontSize=12, spaceAfter=10)))
        
        features = [
            "Multi-tenant Architecture with subdomain-based tenant isolation",
            "Employee Portal for self-service shift viewing and swap requests",
            "Admin Panel with comprehensive roster management and RBAC",
            "Developer Panel for super-admin controls and tenant management",
            "Google Sheets Integration for CSV import/export functionality",
            "Real-time Notifications via email and in-app system"
        ]
        
        for feature in features:
            story.append(Paragraph(f"• {feature}", self.body_style))
        
        story.append(PageBreak())
        
        # Architecture Overview
        story.append(Paragraph("2. System Architecture Overview", self.heading_style))
        story.append(Paragraph(
            "RosterBhai follows a 3-tier architecture with clear separation of concerns:",
            self.body_style
        ))
        story.append(Spacer(1, 0.2*inch))
        
        arch_layers = [
            ["Layer", "Technology", "Purpose"],
            ["Presentation", "Next.js App Router + React", "User Interface"],
            ["Business Logic", "Next.js API Routes", "Server-side Processing"],
            ["Data", "JSON File Storage", "Tenant-isolated Data"]
        ]
        
        table = Table(arch_layers, colWidths=[2*inch, 2.5*inch, 2.5*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 11),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(table)
        story.append(PageBreak())
        
        # Multi-Tenant Architecture
        story.append(Paragraph("3. Multi-Tenant Architecture", self.heading_style))
        story.append(Paragraph(
            "RosterBhai implements subdomain-based multi-tenancy for complete tenant isolation:",
            self.body_style
        ))
        story.append(Spacer(1, 0.2*inch))
        
        story.append(Paragraph("Main Domain: rosterbhai.me", 
                              ParagraphStyle('mono', fontName='Courier', fontSize=10, spaceAfter=5)))
        story.append(Paragraph("Tenant Subdomains: [tenant-slug].rosterbhai.me", 
                              ParagraphStyle('mono', fontName='Courier', fontSize=10, spaceAfter=10)))
        
        story.append(Paragraph(
            "Each tenant has completely isolated data storage in separate directories, ensuring "
            "data privacy and security. The middleware enforces routing rules to prevent "
            "cross-tenant data access.",
            self.body_style
        ))
        story.append(PageBreak())
        
        # System Components
        story.append(Paragraph("4. System Components", self.heading_style))
        
        components = [
            ("Landing Page", "Public-facing marketing website with company registration"),
            ("Developer Portal", "Super-admin panel for SaaS management and tenant approval"),
            ("Admin Panel", "Tenant-level roster management with RBAC support"),
            ("Employee Portal", "Self-service portal for schedule viewing and requests")
        ]
        
        for comp_name, comp_desc in components:
            story.append(Paragraph(f"<b>{comp_name}</b>", 
                                  ParagraphStyle('bold', fontName='Helvetica-Bold', fontSize=12, spaceAfter=5)))
            story.append(Paragraph(comp_desc, self.body_style))
            story.append(Spacer(1, 0.1*inch))
        
        story.append(PageBreak())
        
        # Data Models
        story.append(Paragraph("5. Data Models", self.heading_style))
        story.append(Paragraph(
            "The system uses a file-based JSON storage with the following core entities:",
            self.body_style
        ))
        story.append(Spacer(1, 0.2*inch))
        
        entities = [
            "Tenant - Organization using the service",
            "Employee - Individual workers with schedules",
            "Admin - Tenant-level administrators",
            "Developer - System super-administrators",
            "Schedule Request - Shift change/swap requests",
            "Roster Data - Shift schedules per team"
        ]
        
        for entity in entities:
            story.append(Paragraph(f"• {entity}", self.body_style))
        
        story.append(PageBreak())
        
        # Add diagrams
        story.append(Paragraph("9. System Diagrams", self.heading_style))
        
        diagrams = [
            ("Use Case Diagram", "diagrams/Use_Case_Diagram.png"),
            ("System Architecture", "diagrams/System_Architecture.png"),
            ("Context DFD", "diagrams/Context_DFD.png"),
            ("Level 1 DFD", "diagrams/Level1_DFD.png"),
            ("Entity Relationship Diagram", "diagrams/ER_Diagram.png")
        ]
        
        base_path = "/home/runner/work/rster/rster/Documentation"
        for diagram_name, diagram_path in diagrams:
            story.append(Paragraph(diagram_name, 
                                  ParagraphStyle('diag_title', fontName='Helvetica-Bold', 
                                               fontSize=14, spaceAfter=10, alignment=TA_CENTER)))
            full_path = os.path.join(base_path, diagram_path)
            if os.path.exists(full_path):
                img = Image(full_path, width=6*inch, height=4*inch, kind='proportional')
                story.append(img)
            story.append(PageBreak())
        
        # Build PDF
        doc.build(story)
        print(f"✓ PDF generated: {output_file}")

class PPTXGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, title, subtitle):
        """Add title slide"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
    
    def add_content_slide(self, title, content_points):
        """Add content slide with bullet points"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        
        title_shape.text = title
        tf = body_shape.text_frame
        
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
    
    def add_image_slide(self, title, image_path):
        """Add slide with image"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.CENTER
        
        # Add image
        if os.path.exists(image_path):
            left = Inches(1)
            top = Inches(1.5)
            slide.shapes.add_picture(image_path, left, top, width=Inches(8), height=Inches(5))
    
    def generate_presentation(self, output_file):
        """Generate complete PowerPoint presentation"""
        
        # Title Slide
        self.add_title_slide(
            "RosterBhai",
            f"Multi-Tenant Roster Management System\nVersion 2.0.0\n{datetime.now().strftime('%B %d, %Y')}"
        )
        
        # Overview
        self.add_content_slide(
            "Executive Summary",
            [
                "Modern SaaS application for roster and shift management",
                "Built with Next.js 14 and TypeScript",
                "Multi-tenant architecture with subdomain isolation",
                "Comprehensive features for admins and employees",
                "Google Sheets integration for easy data import",
                "Real-time notifications and request management"
            ]
        )
        
        # Architecture
        self.add_content_slide(
            "System Architecture",
            [
                "3-Tier Architecture Pattern",
                "Presentation: Next.js App Router + React Components",
                "Business Logic: Next.js API Routes (Serverless)",
                "Data: File-based JSON storage with tenant isolation",
                "Middleware-based routing and security",
                "RESTful API design"
            ]
        )
        
        # Multi-Tenancy
        self.add_content_slide(
            "Multi-Tenant Architecture",
            [
                "Subdomain-based tenant identification",
                "Main domain: rosterbhai.me (landing + developer portal)",
                "Tenant subdomains: [slug].rosterbhai.me",
                "Complete data isolation per tenant",
                "Separate data directories for each tenant",
                "Middleware enforces routing rules"
            ]
        )
        
        # Components
        self.add_content_slide(
            "System Components",
            [
                "Landing Page - Marketing and company registration",
                "Developer Portal - Tenant approval and management",
                "Admin Panel - Roster management with RBAC",
                "Employee Portal - Self-service schedule viewing",
                "CSV Import System - Bulk data import",
                "Notification System - In-app and email alerts"
            ]
        )
        
        # User Roles
        self.add_content_slide(
            "User Roles & Capabilities",
            [
                "Developer: Approve tenants, manage subscriptions, CMS",
                "Admin: Manage employees, modify rosters, approve requests",
                "Employee: View schedules, submit change/swap requests",
                "RBAC Support: Role-based permissions for admins",
                "Session-based authentication for all roles",
                "Secure cookie-based session management"
            ]
        )
        
        # Features
        self.add_content_slide(
            "Key Features",
            [
                "CSV & Google Sheets roster import",
                "Shift template management",
                "Employee profile management with photos",
                "Shift change and swap request workflow",
                "Real-time notification system",
                "Audit trail for all modifications",
                "Team-based schedule organization",
                "Mobile-responsive design"
            ]
        )
        
        # Add diagram slides
        base_path = "/home/runner/work/rster/rster/Documentation"
        diagrams = [
            ("Use Case Diagram", "diagrams/Use_Case_Diagram.png"),
            ("System Architecture", "diagrams/System_Architecture.png"),
            ("Context Data Flow Diagram", "diagrams/Context_DFD.png"),
            ("Level 1 Data Flow Diagram", "diagrams/Level1_DFD.png"),
            ("Entity Relationship Diagram", "diagrams/ER_Diagram.png")
        ]
        
        for diagram_title, diagram_path in diagrams:
            full_path = os.path.join(base_path, diagram_path)
            self.add_image_slide(diagram_title, full_path)
        
        # Security
        self.add_content_slide(
            "Security Features",
            [
                "HTTP-only cookies prevent XSS attacks",
                "Subdomain-based tenant isolation",
                "Session validation on every request",
                "Input validation with Zod schemas",
                "File upload restrictions and validation",
                "Role-based access control (RBAC)",
                "Audit logging for accountability"
            ]
        )
        
        # Technology Stack
        self.add_content_slide(
            "Technology Stack",
            [
                "Frontend: Next.js 14, React 18, TypeScript 5",
                "Styling: CSS Modules, Modern UI",
                "Icons: Lucide React",
                "File Processing: Formidable, csv-parse",
                "Deployment: Vercel (recommended)",
                "Storage: JSON files (migrate to DB for scale)",
                "Session: Cookie-based with HTTP-only flag"
            ]
        )
        
        # Future Enhancements
        self.add_content_slide(
            "Future Enhancements",
            [
                "Native mobile apps (iOS/Android)",
                "Home screen widget for shifts",
                "Push notifications",
                "Advanced analytics and reports",
                "SSO (Single Sign-On) integration",
                "Webhook support for third-party integration",
                "Auto-approval rules for requests",
                "Smart scheduling algorithms"
            ]
        )
        
        # Conclusion
        self.add_content_slide(
            "Conclusion",
            [
                "Robust, scalable foundation for roster management",
                "Clear separation of concerns",
                "Multi-tenancy with complete isolation",
                "Comprehensive features for all user roles",
                "Ready for production deployment",
                "Clear upgrade path for scale",
                "Modern technology stack",
                "Focus on security and data privacy"
            ]
        )
        
        # Save presentation
        self.prs.save(output_file)
        print(f"✓ PowerPoint presentation generated: {output_file}")

def main():
    print("Generating PDF and PowerPoint documentation...\n")
    
    # Generate PDF
    pdf_path = "/home/runner/work/rster/rster/Documentation/pdfs/RosterBhai_Complete_Documentation.pdf"
    os.makedirs(os.path.dirname(pdf_path), exist_ok=True)
    pdf_gen = PDFGenerator()
    pdf_gen.generate_comprehensive_pdf(pdf_path)
    
    # Generate PowerPoint
    pptx_path = "/home/runner/work/rster/rster/Documentation/presentations/RosterBhai_Presentation.pptx"
    os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
    pptx_gen = PPTXGenerator()
    pptx_gen.generate_presentation(pptx_path)
    
    print("\n✓ All PDF and PowerPoint files generated successfully!")
    print("\nOutput files:")
    print(f"  - {pdf_path}")
    print(f"  - {pptx_path}")

if __name__ == "__main__":
    main()
